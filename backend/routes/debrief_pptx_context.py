"""Pull recent post-match coach reports (.pptx) from GCS and extract their
text content for prompt injection.

Coaches put their tactical truth in PowerPoint decks under
gs://hfc-football-dev-gcs-coaching-data-dev/. Each deck has a fixed skeleton:
Analytics Summary, Copilot Summary, coach summaries by line, Structure, Oppo,
Development, Review Focuses. We harvest the analyst-written text and feed it
to the model so the AI debrief picks up recurring themes and club jargon.
"""

import io
import logging
import os
import re
from pathlib import Path
from typing import Optional

logger = logging.getLogger(__name__)

REPORTS_BUCKET = os.environ.get(
    "COACH_REPORTS_BUCKET", "hfc-football-dev-gcs-coaching-data-dev"
)
LOCAL_CACHE_DIR = Path("/tmp/hfc-coach-reports-cache")


def _list_pptx_blobs():
    from google.cloud import storage
    client = storage.Client()
    bucket = client.bucket(REPORTS_BUCKET)
    return [b for b in bucket.list_blobs() if b.name.lower().endswith(".pptx")]


def _sort_key(name: str):
    """Reports are named like '4. GEEL.pptx', '5. WB.pptx'. Sort by leading number."""
    m = re.match(r"\s*(\d+)", os.path.basename(name))
    return (int(m.group(1)) if m else 9999, name)


def _extract_pptx_text(local_path: Path) -> str:
    """Pull all text + table content from a .pptx file."""
    from pptx import Presentation
    prs = Presentation(local_path)
    out: list[str] = []
    for i, slide in enumerate(prs.slides, 1):
        slide_lines: list[str] = []
        title = ""
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            txt = "\n".join(p.text for p in shape.text_frame.paragraphs if p.text.strip())
            if not txt.strip():
                continue
            if shape == slide.shapes.title or shape.name.lower().startswith("title"):
                title = txt
            else:
                slide_lines.append(txt)
        # Tables — coach summaries live in tables
        for shape in slide.shapes:
            if shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip().replace("\n", " | ") for c in row.cells if c.text.strip()]
                    if cells:
                        slide_lines.append(" | ".join(cells))
        if not (title or slide_lines):
            continue
        out.append(f"--- Slide {i} ---")
        if title:
            out.append(f"TITLE: {title}")
        out.extend(slide_lines)
    return "\n".join(out)


def _cached_extract(blob) -> Optional[str]:
    """Download blob if not cached, then return its extracted text."""
    LOCAL_CACHE_DIR.mkdir(parents=True, exist_ok=True)
    safe_name = blob.name.replace("/", "_").replace(" ", "_")
    pptx_path = LOCAL_CACHE_DIR / safe_name
    text_path = LOCAL_CACHE_DIR / f"{safe_name}.txt"

    if text_path.exists() and text_path.stat().st_mtime >= blob.updated.timestamp():
        return text_path.read_text(encoding="utf-8")

    try:
        blob.download_to_filename(pptx_path)
        text = _extract_pptx_text(pptx_path)
        text_path.write_text(text, encoding="utf-8")
        return text
    except Exception as e:
        logger.warning(f"Failed to extract {blob.name}: {e}")
        return None


def fetch_recent_reports_context(
    n: int = 3, exclude_round: Optional[str] = None
) -> str:
    """Return a markdown-shaped block summarising the last N coach reports.

    Each report's full text is included — the model can pick out recurring RFIs,
    coaching jargon, and review themes. ~5-10K tokens for n=3.

    Args:
        n: how many recent reports to include
        exclude_round: e.g. 'R6' or '6' — skip the report for the round being
                       generated to avoid the AI just regurgitating it.
    """
    try:
        blobs = _list_pptx_blobs()
    except Exception as e:
        logger.warning(f"Could not list reports bucket: {e}")
        return ""
    if not blobs:
        return ""

    blobs.sort(key=lambda b: _sort_key(b.name), reverse=True)

    out: list[str] = []
    used = 0
    for blob in blobs:
        if used >= n:
            break
        if exclude_round:
            base = os.path.basename(blob.name)
            num = re.match(r"\s*(\d+)", base)
            if num and num.group(1) == str(exclude_round).lstrip("R").lstrip("r"):
                continue
        text = _cached_extract(blob)
        if not text:
            continue
        out.append(f"\n### Coach Report — {os.path.basename(blob.name)}\n{text}\n")
        used += 1

    if not out:
        return ""

    return (
        "PRIOR COACH REPORTS (for context — recurring RFIs, club jargon, "
        "review themes; do not just repeat them, build on them):\n"
        + "\n".join(out)
    )
